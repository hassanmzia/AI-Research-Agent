"""Generate Technical Architecture PowerPoint presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_BLUE = RGBColor(0x16, 0x21, 0x3e)
CARD_BG = RGBColor(0x0f, 0x3d, 0x6e)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_CYAN = RGBColor(0x00, 0xBC, 0xD4)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_PURPLE = RGBColor(0x9C, 0x27, 0xB0)
ACCENT_RED = RGBColor(0xE9, 0x1E, 0x63)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill_color, text="", font_size=10, bold=False, align=PP_ALIGN.LEFT, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_WHITE
        p.font.bold = bold
        p.alignment = align
    return shape

def add_text(slide, left, top, width, height, text, font_size=12, bold=False, color=TEXT_WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_arrow(slide, x1, y1, x2, y2, color=ACCENT_CYAN):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)

# ==================== SLIDE 1: Title ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide1)

add_text(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "AI Research Multi-Agent System", font_size=44, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
add_text(slide1, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
         "Technical Architecture Overview", font_size=24, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# Subtitle cards
for i, (label, val) in enumerate([
    ("Multi-Agent Pipeline", "5 Phases"),
    ("AGI Evaluation", "10 Parameters"),
    ("Docker Services", "7 Containers"),
    ("API Endpoints", "30+"),
    ("Real-time Updates", "WebSocket"),
]):
    x = Inches(1.5 + i * 2.1)
    add_box(slide1, x, Inches(4.2), Inches(1.9), Inches(1.0), CARD_BG, border_color=ACCENT_CYAN)
    add_text(slide1, x + Pt(8), Inches(4.3), Inches(1.7), Inches(0.4), val, font_size=18, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
    add_text(slide1, x + Pt(8), Inches(4.7), Inches(1.7), Inches(0.4), label, font_size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_text(slide1, Inches(1), Inches(6), Inches(11), Inches(0.5),
         "Django + React + Node.js + PostgreSQL + Redis + Celery + LangChain + OpenAI",
         font_size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# ==================== SLIDE 2: System Architecture ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

add_text(slide2, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
         "System Architecture", font_size=28, bold=True, color=TEXT_WHITE)

# Layer labels and boxes
layers = [
    ("USER LAYER", Inches(0.9), RGBColor(0x15, 0x65, 0xC0), [
        ("React 18 + TS", "Tailwind CSS\nZustand State"),
        ("10 Pages", "Dashboard, Sessions\nPapers, Collections"),
        ("Nginx", ":3045\nStatic + Proxy"),
    ]),
    ("API GATEWAY", Inches(2.1), RGBColor(0xF5, 0x7C, 0x00), [
        ("Express.js", ":3046\nJWT | Rate Limit"),
        ("HTTP Proxy", "/api/* -> backend\nPath Rewriting"),
        ("WebSocket", "/ws Real-time\nSession Updates"),
        ("Redis Sub", "Pub/Sub Channel\nresearch_updates"),
    ]),
    ("BACKEND", Inches(3.3), RGBColor(0x2E, 0x7D, 0x32), [
        ("Django 5.1 + DRF", ":8045\nGunicorn 4 workers"),
        ("REST API", "30+ Endpoints\nCRUD + Actions"),
        ("MCP Server", "JSON-RPC 2.0\n5 Tools"),
        ("A2A Protocol", "Agent Cards\nTask Management"),
    ]),
    ("MULTI-AGENT PIPELINE", Inches(4.5), RGBColor(0x9C, 0x27, 0xB0), [
        ("Lead Supervisor", "Orchestrator\nState Machine"),
        ("Planner Agent", "GPT-4o-mini\nQuery Generation"),
        ("Discovery Agent", "ArXiv API\nRelevance Filter"),
        ("Evaluation Agent", "10-Param AGI\nScoring"),
        ("Synthesis Agent", "Report Gen\nMarkdown"),
    ]),
    ("ASYNC PROCESSING", Inches(5.7), RGBColor(0xE6, 0x51, 0x00), [
        ("Celery Worker", "Pipeline Tasks\nAsync Execution"),
        ("Celery Beat", "Scheduled Jobs\nDaily/Weekly"),
        ("Redis Publisher", "Phase Changes\nNotifications"),
        ("OpenAI API", "GPT-4o-mini\nLLM Calls"),
        ("ArXiv API", "Paper Search\nexport.arxiv.org"),
    ]),
    ("DATA LAYER", Inches(6.7), RGBColor(0x19, 0x76, 0xD2), [
        ("PostgreSQL 16", ":5445\n8 Tables"),
        ("Redis 7", ":6345\nQueue + Cache"),
        ("Docker Volumes", "pgdata\nstaticfiles"),
    ]),
]

for label, y, color, items in layers:
    # Layer label
    add_box(slide2, Inches(0.3), y, Inches(1.5), Inches(0.9), color, label, font_size=9, bold=True, align=PP_ALIGN.CENTER)
    # Items
    item_w = Inches((10.5) / len(items))
    for i, (title, desc) in enumerate(items):
        x = Inches(2.0 + i * (10.5 / len(items)))
        box = add_box(slide2, x, y, item_w - Inches(0.1), Inches(0.9), RGBColor(0x22, 0x33, 0x55), border_color=color)
        tf = box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(8)
        p2.font.color.rgb = TEXT_LIGHT

# ==================== SLIDE 3: Agent Pipeline ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)

add_text(slide3, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
         "Multi-Agent Research Pipeline", font_size=28, bold=True, color=TEXT_WHITE)

# Pipeline phases
phases = [
    ("1. PLANNING", "Planner Agent", "GPT-4o-mini generates:\n- Search keywords\n- ArXiv queries (AND logic)\n- Required terms\n- Category selection\n- Focus areas", RGBColor(0x15, 0x65, 0xC0)),
    ("2. DISCOVERY", "Discovery Agent", "ArXiv API search:\n- Multiple targeted queries\n- Relevance sort\n- Progressive broadening\n- Deduplication\n- Relevance scoring & filter", RGBColor(0x2E, 0x7D, 0x32)),
    ("3. EVALUATION", "Evaluation Agent", "10-Parameter AGI scoring:\n- LLM scores each paper\n- Weighted average (1-100)\n- Classification: High/Med/Low\n- Key innovations\n- Limitations analysis", RGBColor(0xE6, 0x51, 0x00)),
    ("4. SYNTHESIS", "Synthesis Agent", "Report generation:\n- Executive summary\n- Paper rankings\n- Score distribution\n- Methodology docs\n- Markdown export", RGBColor(0x88, 0x0E, 0x4F)),
]

for i, (phase, agent, desc, color) in enumerate(phases):
    x = Inches(0.5 + i * 3.2)
    # Phase header
    add_box(slide3, x, Inches(1.0), Inches(3.0), Inches(0.5), color, phase, font_size=14, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide3, x, Inches(1.5), Inches(3.0), Inches(0.35), RGBColor(0x22, 0x33, 0x55), agent, font_size=11, bold=True, align=PP_ALIGN.CENTER, border_color=color)
    # Description
    box = add_box(slide3, x, Inches(1.9), Inches(3.0), Inches(2.5), RGBColor(0x1a, 0x25, 0x40), border_color=color)
    tf = box.text_frame
    tf.paragraphs[0].text = ""
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(2)

# Arrow connectors between phases
for i in range(3):
    x1 = Inches(3.5 + i * 3.2)
    x2 = x1 + Inches(0.2)
    y = Inches(2.5)
    add_text(slide3, x1 - Inches(0.05), Inches(2.2), Inches(0.3), Inches(0.3), "->", font_size=20, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# AGI Parameters section
add_text(slide3, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
         "AGI Evaluation Parameters (Weighted)", font_size=16, bold=True, color=ACCENT_CYAN)

params = [
    ("Novel Problem Solving", "15%"), ("Few-Shot Learning", "15%"), ("Task Transfer", "15%"),
    ("Abstract Reasoning", "12%"), ("Contextual Adaptation", "10%"), ("Multi-Rule Integration", "10%"),
    ("Generalization Efficiency", "8%"), ("Meta-Learning", "8%"), ("World Modeling", "4%"),
    ("Autonomous Goal Setting", "3%"),
]
for i, (name, weight) in enumerate(params):
    row = i // 5
    col = i % 5
    x = Inches(0.5 + col * 2.5)
    y = Inches(5.1 + row * 0.55)
    colors = [ACCENT_GREEN, ACCENT_GREEN, ACCENT_GREEN, ACCENT_CYAN, ACCENT_ORANGE, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_PURPLE, TEXT_LIGHT, TEXT_LIGHT]
    add_box(slide3, x, y, Inches(2.3), Inches(0.45), RGBColor(0x22, 0x33, 0x55), f"{name} ({weight})", font_size=9, bold=False, align=PP_ALIGN.CENTER, border_color=colors[i])

# Score interpretation
add_text(slide3, Inches(0.5), Inches(6.3), Inches(12), Inches(0.3),
         "Score: 90-100 Exceptional  |  70-89 High AGI  |  40-69 Medium  |  0-39 Low",
         font_size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# ==================== SLIDE 4: Tech Stack ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)

add_text(slide4, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
         "Technology Stack & Infrastructure", font_size=28, bold=True, color=TEXT_WHITE)

stacks = [
    ("Frontend", ACCENT_CYAN, [
        "React 18 + TypeScript 5.7",
        "Vite 6.0 (Build)",
        "Tailwind CSS 3.4",
        "Zustand 5.0 (State)",
        "Recharts (Visualization)",
        "Axios (HTTP Client)",
        "WebSocket (Real-time)",
        "Lucide React (Icons)",
    ]),
    ("API Gateway", ACCENT_ORANGE, [
        "Node.js + Express 4.21",
        "http-proxy-middleware",
        "ws (WebSocket Server)",
        "ioredis (Pub/Sub)",
        "jsonwebtoken (JWT)",
        "helmet (Security)",
        "express-rate-limit",
        "compression + morgan",
    ]),
    ("Backend", ACCENT_GREEN, [
        "Django 5.1 + DRF 3.15",
        "SimpleJWT (Auth)",
        "Celery 5.4 (Async)",
        "Celery Beat (Scheduler)",
        "LangChain 0.3.14",
        "LangGraph 0.2.60",
        "drf-spectacular (Docs)",
        "Gunicorn (WSGI)",
    ]),
    ("AI / ML", ACCENT_PURPLE, [
        "OpenAI GPT-4o-mini",
        "LangChain + LangGraph",
        "arxiv 2.1.3 (Discovery)",
        "tenacity (Retry Logic)",
        "MCP Protocol 1.2.0",
        "A2A Protocol (Google)",
        "10-Param AGI Framework",
        "Relevance Scoring Engine",
    ]),
    ("Infrastructure", ACCENT_RED, [
        "Docker + Compose 3.9",
        "PostgreSQL 16-alpine",
        "Redis 7-alpine",
        "Nginx (Reverse Proxy)",
        "7 Docker Containers",
        "Bridge Network",
        "Health Checks",
        "Volume Persistence",
    ]),
]

for i, (title, color, items) in enumerate(stacks):
    x = Inches(0.3 + i * 2.6)
    add_box(slide4, x, Inches(0.9), Inches(2.4), Inches(0.45), color, title, font_size=14, bold=True, align=PP_ALIGN.CENTER)
    box = add_box(slide4, x, Inches(1.4), Inches(2.4), Inches(4.2), RGBColor(0x1a, 0x25, 0x40), border_color=color)
    tf = box.text_frame
    tf.paragraphs[0].text = ""
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(6)

# Docker ports table
add_text(slide4, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
         "Port Mapping:  Frontend :3045  |  Gateway :3046  |  Backend :8045  |  PostgreSQL :5445  |  Redis :6345",
         font_size=12, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# ==================== SLIDE 5: Data Flow ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)

add_text(slide5, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
         "Request Flow: User to Research Results", font_size=28, bold=True, color=TEXT_WHITE)

steps = [
    ("1", "User submits\nresearch objective", ACCENT_CYAN),
    ("2", "React app sends\nPOST /api/sessions/", ACCENT_CYAN),
    ("3", "Gateway validates JWT\nproxies to backend", ACCENT_ORANGE),
    ("4", "Django creates session\nlaunches Celery task", ACCENT_GREEN),
    ("5", "Celery worker runs\nmulti-agent pipeline", ACCENT_PURPLE),
    ("6", "Pipeline phases:\nPlan -> Discover ->\nEvaluate -> Synthesize", ACCENT_PURPLE),
    ("7", "Redis publishes\nphase updates", ACCENT_RED),
    ("8", "Gateway broadcasts\nvia WebSocket", ACCENT_ORANGE),
    ("9", "Frontend updates\nin real-time", ACCENT_CYAN),
    ("10", "Results saved to\nPostgreSQL", RGBColor(0x19, 0x76, 0xD2)),
]

for i, (num, desc, color) in enumerate(steps):
    row = i // 5
    col = i % 5
    x = Inches(0.5 + col * 2.5)
    y = Inches(1.0 + row * 2.8)

    # Number circle
    circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEXT_WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Description
    add_box(slide5, x, y + Inches(0.6), Inches(2.2), Inches(1.5), RGBColor(0x1a, 0x25, 0x40), desc, font_size=11, align=PP_ALIGN.CENTER, border_color=color)

# ==================== SLIDE 6: Key Features ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)

add_text(slide6, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
         "Key Features & Capabilities", font_size=28, bold=True, color=TEXT_WHITE)

features = [
    ("Intelligent Paper Discovery", "ArXiv API with targeted AND queries, relevance sorting, progressive broadening, and post-fetch relevance filtering with required term gating", ACCENT_GREEN),
    ("10-Parameter AGI Evaluation", "Comprehensive scoring framework: Novel Problem Solving, Few-Shot Learning, Task Transfer, Abstract Reasoning + 6 more weighted parameters", ACCENT_PURPLE),
    ("Real-time Pipeline Updates", "WebSocket-based live updates through Redis pub/sub. Users see each phase complete in real-time with agent activity logs", ACCENT_CYAN),
    ("MCP + A2A Protocol Support", "Model Context Protocol (JSON-RPC 2.0) with 5 tools and Google A2A agent-to-agent communication with agent cards", ACCENT_ORANGE),
    ("Collections & Organization", "User-curated paper collections across sessions. Add/remove papers, browse collection contents, public/private visibility", RGBColor(0x19, 0x76, 0xD2)),
    ("Scheduled Research", "Automated recurring research: daily, weekly, biweekly, monthly. Celery Beat scheduler with run tracking and notifications", ACCENT_RED),
    ("Multi-Format Export", "Export research reports as Markdown, JSON, CSV, PDF, or Excel. Full data preservation including evaluations and scores", ACCENT_TEAL),
    ("Interactive Dashboard", "Analytics with KPI cards, score distribution charts, papers by source, sessions over time, and top-performing papers", RGBColor(0xFF, 0xD5, 0x4F)),
]

for i, (title, desc, color) in enumerate(features):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(0.9 + row * 1.55)

    add_box(slide6, x, y, Inches(6.1), Inches(0.35), color, title, font_size=13, bold=True, align=PP_ALIGN.LEFT)
    add_box(slide6, x, y + Inches(0.38), Inches(6.1), Inches(1.0), RGBColor(0x1a, 0x25, 0x40), desc, font_size=10, border_color=color)

# Save
output_path = "/home/user/AI-Research-Agent/AI_Research_Agent_Architecture.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
